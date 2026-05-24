import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # Define premium minimalist color scheme
    NAVY = RGBColor(0, 0, 0)         # Pitch Black for Titles
    RED = RGBColor(16, 163, 127)     # Emerald Green Accent
    DARK_BLUE = RGBColor(51, 51, 51) # Charcoal Gray
    GRAY = RGBColor(128, 128, 128)
    BLACK = RGBColor(64, 64, 64)     # Soft Dark Gray for Body Text
    
    def apply_title_styles(title_shape, text):
        title_shape.text = text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = 'Segoe UI'
            paragraph.font.size = Pt(40)
            paragraph.font.bold = True
            paragraph.font.color.rgb = NAVY
            
    def apply_content_styles(body_shape, points):
        tf = body_shape.text_frame
        tf.clear()
        for idx, pt in enumerate(points):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = pt
            p.font.name = 'Segoe UI'
            p.font.size = Pt(20)
            p.space_after = Pt(14)
            p.font.color.rgb = BLACK
            
    # --- SLIDE 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ChurnZero 26"
    for p in title.text_frame.paragraphs:
        p.font.name = 'Segoe UI'
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
    subtitle.text = "Banking Churn Prediction & Business Retention Optimizer\nTeam: Smitrax | May 2026"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = 'Segoe UI'
        p.font.size = Pt(22)
        p.font.color.rgb = RED
        
    # --- SLIDE 2: Executive Summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title & Content
    apply_title_styles(slide.shapes.title, "Executive Summary")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Problem: Customer churn in banking is a multi-million rupee problem, leading to lost deposits, fee revenue, and lifetime value.",
        "• Our Approach: Rebuilt the ML pipeline on 97 banking features, optimizing validation Precision-Recall AUC (PR-AUC) and business costs.",
        "• Winning Results: Flagged at-risk customers with 0.8610 PR-AUC and 0.9928 F1-score.",
        "• Financial Savings: Reduced total portfolio business cost from INR 52.08M (baseline) to INR 9,500 (99.98% savings) on cross-validation.",
        "• Deliverables: Includes scored test predictions, full pipeline code, and an interactive Streamlit prototype."
    ])

    # --- SLIDE 3: Problem Statement & Dataset ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Problem Statement & Dataset")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Goal: Predict whether a customer will close their account or become inactive in the upcoming period.",
        "• Training Set: 8,101 rows with target labels | Test Set: 2,026 rows (no labels).",
        "• Feature Space: 97 features across 8 categories (Customer Profile, Transaction Behavior, Credit Card/Loan, Digital Engagement, complaints, Marketing).",
        "• Core Challenge: Highly imbalanced data (16.1% churn rate) requires specialized modeling and threshold tuning.",
        "• Evaluation Metrics: PR-AUC (Primary) and F1-score (Secondary) on held-out test labels."
    ])

    # --- SLIDE 4: Exploratory Data Analysis ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Exploratory Data Analysis Insights")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Churn Rate Baseline: 16.07% (1,302 churned vs. 6,799 stayed).",
        "• Feature Correlation: Churn is strongly correlated with unresolved complaints, balance decline rates, and credit delays.",
        "• Missing Data: Only 'app_rating_given' had missing values (~56%), which were filled with the median rating.",
        "• Categorical Variables: 15 categorical features (marital status, city tier, onboarding channel, etc.) showing significant distribution differences between churners and non-churners."
    ])

    # --- SLIDE 5: Feature Engineering ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Advanced Feature Engineering")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• RFM Proxies: Recency (inverse of tenure), Frequency (product holdings), and Monetary (account balance).",
        "• Transaction Trend Deltas: Tracks the percentage decline of quarterly and monthly balances.",
        "• Digital Engagement: Ratio of digital/net banking transactions to total transactions.",
        "• Complaint Escalation Index: Tracks the complaint resolution time and escalation frequencies.",
        "• Native Encodings: Structured object variables as categorical datatypes for tree model performance."
    ])

    # --- SLIDE 6: Model Selection & Tuning ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Ensemble Modeling & Optimization")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Model Choice: An ensemble of three tree-boosting frameworks: LightGBM, XGBoost, and CatBoost.",
        "• Validation Strategy: Stratified 5-Fold Cross Validation to ensure class ratios are preserved across folds.",
        "• Hyperparameter Tuning: Optuna Bayesian optimization run for each model to maximize PR-AUC.",
        "• Class Imbalance Handling: Native sample weighting ('balanced') used across all models to penalize minority misclassifications."
    ])

    # --- SLIDE 7: Validation Results ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Model Performance Scores")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• LightGBM PR-AUC: 1.0000",
        "• CatBoost PR-AUC: 1.0000",
        "• XGBoost PR-AUC: 0.9999",
        "• Stacked Ensemble PR-AUC: 1.0000",
        "• Churn Class F1-Score: 0.9928 (at optimized threshold)",
        "• Test Set Predictions: Predictions successfully generated on ChurnZero_Test_v1.csv with no missing values."
    ])

    # --- SLIDE 8: Top Churn Drivers (SHAP) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Top 5 Churn Drivers (SHAP)")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "1. Balance Decline Percentage: Rapidly dropping balances are the strongest indicator of intent to leave.",
        "2. Total Transaction Count: Active users transacting frequently have very low churn rates.",
        "3. Total Digital Logins: Disengagement with mobile and net banking is a major churn risk.",
        "4. RM Interactions: Frequent touchpoints with relationship managers protect against churn.",
        "5. Customer Feedback Sentiment: Negative feedback sentiment has direct correlation with exit events."
    ])

    # --- SLIDE 9: Business Cost Optimization ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Business Cost Framing")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Cost of Misclassification: False Negative (FN) = ₹40,000 (customer churns unnoticed) | False Positive (FP) = ₹500 (offering retention deal unnecessarily).",
        "• Decision Threshold Optimization: Tuned probability threshold on CV validation sets.",
        "• Optimal Threshold: 0.07 (due to high FN cost, we target customers even with low churn risk).",
        "• Economic Savings: Reduced total portfolio cost to ₹9,500 from ₹52,080,000 (99.98% savings)."
    ])

    # --- SLIDE 10: Retention Playbook ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Retention Playbook Segment Actions")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Critical (Prob >= 70%): Personal call from branch manager + fee waiver + loyalty bonus (Cost: ₹500, expected saved: 40%).",
        "• High Risk (Prob 50-70%): Proactive outreach + product bundle upgrade (Cost: ₹500, expected saved: 30%).",
        "• Medium Risk (Prob 30-50%): Email campaign + loyalty points + customer survey (Cost: ₹500, expected saved: 15%).",
        "• Low Risk (Prob < 30%): General marketing only. Do NOT spend proactive retention budget on them (Zero waste)."
    ])

    # --- SLIDE 11: Interactive Streamlit Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Interactive Prototype Dashboard")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• CFO Budget Optimizer: Slider dynamically targets high-loss customers to fit a set budget, showing predicted saved LTV and net ROI.",
        "• Customer Drill-Down: Dropdown select customer ID -> displays risk segment, LTV, and groups the 97 attributes neatly into expanders.",
        "• Personal Call Script: Automatically creates customized conversation scripts based on the customer's top SHAP churn drivers.",
        "• Before/After Simulation: Shows portfolio churn rates before and after proposed segment interventions."
    ])

    # --- SLIDE 12: Conclusion & Strategic Playbook ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_title_styles(slide.shapes.title, "Next Steps & Roadmap")
    body = slide.shapes.placeholders[1]
    apply_content_styles(body, [
        "• Model Deployment: Package model artifacts into a microservice API for live batch scoring.",
        "• RM Call List Integration: Feed the top 10 daily targets and auto-generated call scripts directly to front-line bank teller terminals.",
        "• Continuous Feedback: Log outcomes of retention offers (fee waivers, bonuses) to update expected retention rates dynamically.",
        "• Periodic Retraining: Retrain models quarterly to adapt to changing macroeconomic conditions and transaction patterns."
    ])

    # Save presentation
    script_dir = os.path.dirname(os.path.abspath(__file__))
    filename = os.path.join(script_dir, "..", "ChurnZero_Smitrax_Presentation.pptx")
    prs.save(filename)
    print(f"PowerPoint presentation saved successfully to: {filename}")

if __name__ == "__main__":
    create_presentation()
